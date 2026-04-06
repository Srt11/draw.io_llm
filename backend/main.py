from fastapi import FastAPI, UploadFile, File, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import JSONResponse
from pydantic import BaseModel
from dotenv import load_dotenv
import anthropic
import fitz  # PyMuPDF
from pptx import Presentation
import io
import re
import os

load_dotenv()

app = FastAPI(title="DrawIO Chatbot API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:3000", "http://localhost:5173"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

client = anthropic.Anthropic(api_key=os.environ.get("ANTHROPIC_API_KEY"))

SYSTEM_PROMPT = """You are a diagram generation assistant that creates draw.io XML diagrams.

When generating a diagram, ALWAYS return ONLY valid draw.io XML. No explanation text, no markdown, no code fences.

The format must be exactly:
<mxGraphModel>
  <root>
    <mxCell id="0"/>
    <mxCell id="1" parent="0"/>
    <!-- nodes and edges here -->
  </root>
</mxGraphModel>

CRITICAL XML RULES - NEVER BREAK THESE:
1. NEVER use bare & in any attribute value. Always use &amp; instead. Example: "Read &amp; Write" not "Read & Write"
2. NEVER use < or > inside attribute values. Use &lt; and &gt;
3. NEVER use unescaped quotes inside quoted attributes
4. All attribute values must be in double quotes
5. Every opened tag must be closed

Rules for diagram layout:
- Every mxCell must have a unique numeric id starting from 2
- Vertex cells need: vertex="1" parent="1" and an mxGeometry child
- Edge cells need: edge="1" source="ID" target="ID" parent="1"
- Space nodes at least 120px apart vertically for flowcharts
- Use style="rounded=1;whiteSpace=wrap;html=1;" for regular boxes
- Use style="rhombus;whiteSpace=wrap;html=1;" for decision diamonds
- Use style="ellipse;whiteSpace=wrap;html=1;" for start/end ovals

When EDITING an existing diagram:
- Keep all existing cell IDs unchanged
- Add new cells with new unique IDs
- Return the complete updated mxGraphModel XML
"""


class ChatRequest(BaseModel):
    message: str
    current_xml: str = ""
    history: list = []


def extract_text_from_pdf(file_bytes: bytes) -> str:
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
    doc.close()
    return text[:8000]


def extract_text_from_pptx(file_bytes: bytes) -> str:
    prs = Presentation(io.BytesIO(file_bytes))
    text = ""
    for i, slide in enumerate(prs.slides):
        text += f"\n--- Slide {i+1} ---\n"
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text += shape.text + "\n"
    return text[:8000]


def sanitize_xml(xml: str) -> str:
    """Fix bare & in XML attribute values that Claude sometimes generates."""
    # Replace bare & not already part of a valid XML entity
    xml = re.sub(r'&(?!(amp|lt|gt|quot|apos|#\d+|#x[0-9a-fA-F]+);)', '&amp;', xml)
    return xml


def extract_diagram_xml(response_text: str) -> tuple[str, str]:
    """Extract and sanitize mxGraphModel XML from Claude's response."""
    # Primary: find mxGraphModel block directly
    graph_match = re.search(r'(<mxGraphModel[\s\S]*?</mxGraphModel>)', response_text)
    if graph_match:
        return sanitize_xml(graph_match.group(1)), ""

    # Fallback: unwrap <diagram> tags
    diagram_match = re.search(r'<diagram>([\s\S]*?)</diagram>', response_text)
    if diagram_match:
        inner = diagram_match.group(1).strip()
        if "<mxGraphModel" in inner:
            gm = re.search(r'(<mxGraphModel[\s\S]*?</mxGraphModel>)', inner)
            if gm:
                return sanitize_xml(gm.group(1)), ""
        xml = f"<mxGraphModel><root><mxCell id=\"0\"/><mxCell id=\"1\" parent=\"0\"/>{inner}</root></mxGraphModel>"
        return sanitize_xml(xml), ""

    return "", response_text


@app.post("/api/chat")
async def chat(request: ChatRequest):
    messages = []
    for msg in request.history[-6:]:
        messages.append({"role": msg["role"], "content": msg["content"]})

    user_content = request.message
    if request.current_xml:
        user_content = f"Current diagram XML:\n{request.current_xml}\n\nUser request: {request.message}"

    messages.append({"role": "user", "content": user_content})

    response = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=4000,
        system=SYSTEM_PROMPT,
        messages=messages
    )

    response_text = response.content[0].text
    xml, chat_text = extract_diagram_xml(response_text)

    return JSONResponse({
        "xml": xml,
        "message": chat_text if chat_text else ("Diagram generated!" if xml else response_text),
        "has_diagram": bool(xml)
    })


@app.post("/api/upload")
async def upload_file(file: UploadFile = File(...)):
    file_bytes = await file.read()
    filename = file.filename.lower()

    if filename.endswith(".pdf"):
        text = extract_text_from_pdf(file_bytes)
        file_type = "PDF"
    elif filename.endswith((".pptx", ".ppt")):
        text = extract_text_from_pptx(file_bytes)
        file_type = "PowerPoint"
    else:
        raise HTTPException(status_code=400, detail="Only PDF and PPTX files are supported")

    if not text.strip():
        raise HTTPException(status_code=400, detail="Could not extract text from file")

    prompt = f"""Analyze this {file_type} content and create a draw.io diagram capturing key concepts, processes, or structure.

Content:
{text}

Return only the mxGraphModel XML, no explanation."""

    response = client.messages.create(
        model="claude-sonnet-4-20250514",
        max_tokens=4000,
        system=SYSTEM_PROMPT,
        messages=[{"role": "user", "content": prompt}]
    )

    response_text = response.content[0].text
    xml, _ = extract_diagram_xml(response_text)

    if not xml:
        raise HTTPException(status_code=500, detail="Failed to generate diagram from file")

    return JSONResponse({
        "xml": xml,
        "message": f"Generated diagram from {file.filename}",
        "filename": file.filename
    })


@app.get("/api/health")
async def health():
    return {"status": "ok"}
